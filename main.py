import os
import sys
from pptx import Presentation
from PIL import Image
import io
import google.generativeai as genai
from dotenv import load_dotenv

def extract_content_from_pptx(pptx_path):
    """
    Extracts text and images from each slide of a PowerPoint presentation.
    Returns a list of dictionaries, where each dict represents a slide.
    """

    if not os.path.exists(pptx_path):
        print(f"Error: File not found at {pptx_path}")
        return None

    prs = Presentation(pptx_path)
    slide_data = []

    for i, slide in enumerate(prs.slides):
        slide_number = i + 1
        print(f"Processing Slide {slide_number}...")
        
        # Extract text from all shapes
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
        
        # Extract images and store them as PIL Image objects
        image_content = []
        for shape in slide.shapes:
            if shape.shape_type == 13: # 13 for picture
                image = shape.image
                image_bytes = image.blob
                try:
                    pil_image = Image.open(io.BytesIO(image_bytes))
                    image_content.append(pil_image)
                except Exception as e:
                    print(f"Warning: Could not process an image on slide {slide_number}: {e}")

        slide_data.append({
            "slide_number": slide_number,
            "text": "\n".join(text_content),
            "images": image_content
        })
        
    return slide_data



load_dotenv() # Loads the GOOGLE_API_KEY from the .env file

def analyze_presentation_with_gemini(all_slide_data):
    """
    Sends all slide data to Gemini and asks it to find inconsistencies.
    """
    API_KEY = os.getenv("GOOGLE_API_KEY")
    if not API_KEY:
        raise ValueError("GOOGLE_API_KEY not found. Please set it in your .env file.")
        
    genai.configure(api_key=API_KEY)

    # --- A General-Purpose System Prompt ---
    system_prompt = """
    You are a highly logical and detail-oriented analysis agent. Your function is to meticulously review the content of any multi-slide presentation and identify factual, logical, or internal inconsistencies.

    Your instructions are to perform a rigorous, slide-by-slide comparison and identify any contradictions or errors, including but not limited to:

    **1. Factual & Numerical Inconsistencies:**
    - **Conflicting Data:** A number, date, name, or statistic on one slide contradicts information on another slide. (e.g., Slide 2 states an event happened in 1995, but Slide 10 refers to the same event as happening in 1996).
    - **Calculation Errors:** Numbers in a table do not sum up correctly, or percentages presented do not add up to 100%.
    - **Inconsistent Terminology:** The same concept or item is referred to by different names across slides, causing confusion (e.g., "Subject A" on one slide is called "Test Group Alpha" on another).

    **2. Logical & Thematic Inconsistencies:**
    - **Contradictory Claims:** A statement or argument on one slide is logically at odds with a statement on another. (e.g., Slide 3 claims "solar power is ineffective in this region," while Slide 8 proposes a solar-powered solution).
    - **Timeline Mismatches:** The sequence of events described or shown in a timeline is not logical or conflicts with dates mentioned elsewhere.
    - **Causal Errors:** A cause-and-effect relationship is described illogically (e.g., claiming an event from 2020 caused an outcome in 2018).

    **Output Format:**
    Provide your findings as a clear, structured list. For each issue found, you MUST provide:
    - The slide numbers involved.
    - The category of the inconsistency (e.g., "Factual Contradiction", "Logical Mismatch", "Timeline Mismatches", "Calculation Errors").
    - A concise description of the problem.

    If you find no inconsistencies after a thorough review, you must state: "No significant factual or logical inconsistencies were found."
    """

    model = genai.GenerativeModel('gemini-2.5-flash', system_instruction=system_prompt)

    user_prompt = ["Analyze the following presentation content for inconsistencies based on your instructions."]

    for slide in all_slide_data:
        # Add a header for each slide's content
        user_prompt.append(f"--- SLIDE {slide['slide_number']} ---")
        
        # Add the text from the slide
        if slide['text']:
            user_prompt.append("TEXT:")
            user_prompt.append(slide['text'])
        
        # Add the images from the slide
        if slide['images']:
            user_prompt.append("IMAGES:")
            for img in slide['images']:
                user_prompt.append(img) # The API can handle PIL images directly
    
    # Generate the content
    try:
        response = model.generate_content(user_prompt)
        return response.text
    except Exception as e:
        return f"An error occurred while communicating with the Gemini API: {e}"

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python main.py <path_to_pptx_file>")
    else:
        filepath = sys.argv[1]
        extracted_data = extract_content_from_pptx(filepath)
        if extracted_data:
            print("\n Extraction Complete.")
            analysis_result = analyze_presentation_with_gemini(extracted_data)
            
            print("\n--- 🔍 AI Analysis Report ---")
            print(analysis_result)
            print("--- End of Report ---")
